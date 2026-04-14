"""
Hermes Agent 源码深度分析 - PowerPoint 生成脚本
基于 hermes-source-analysis 目录下的 Markdown 文档生成技术分析演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色方案 ──
BG_PRIMARY    = RGBColor(0x0F, 0x17, 0x2A)   # 深海军蓝
BG_SECONDARY  = RGBColor(0x1E, 0x29, 0x3B)   # 暗石板
BG_CARD       = RGBColor(0x27, 0x33, 0x4D)   # 卡片背景
ACCENT_BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # 亮蓝
ACCENT_LIGHT  = RGBColor(0x60, 0xA5, 0xFA)   # 浅蓝
ACCENT_CYAN   = RGBColor(0x22, 0xD3, 0xEE)   # 青色
ACCENT_GREEN  = RGBColor(0x10, 0xB9, 0x81)   # 翠绿
ACCENT_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # 琥珀
ACCENT_ROSE   = RGBColor(0xF4, 0x72, 0xB6)   # 玫粉
ACCENT_RED    = RGBColor(0xEF, 0x44, 0x44)   # 红色
TEXT_WHITE    = RGBColor(0xF1, 0xF5, 0xF9)   # 白色文字
TEXT_LIGHT    = RGBColor(0xCB, 0xD5, 0xE1)   # 浅灰文字
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # 弱化文字
BORDER_COLOR  = RGBColor(0x33, 0x41, 0x55)   # 边框

# ── 字体设置 ──
FONT_CN = "微软雅黑"
FONT_EN = "Calibri"

# ── 演示文稿初始化 ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# 使用空白版式
blank_layout = prs.slide_layouts[6]


def add_bg(slide, color=BG_PRIMARY):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """添加矩形形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_rounded_shape(slide, left, top, width, height, fill_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_CN):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=14,
                       default_color=TEXT_LIGHT, line_spacing=1.3):
    """添加多行文本（支持逐行设置格式）"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, size, color, bold = line_info, default_size, default_color, False
        else:
            text = line_info.get("text", "")
            size = line_info.get("size", default_size)
            color = line_info.get("color", default_color)
            bold = line_info.get("bold", False)

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = FONT_CN
        p.space_after = Pt(4)
        p.line_spacing = Pt(size * line_spacing)

    return txBox


def add_header_bar(slide, title, subtitle=None):
    """添加顶部标题栏"""
    # 蓝色顶部条
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), fill_color=ACCENT_BLUE)
    # 标题文字
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.7),
                 title, font_size=28, color=TEXT_WHITE, bold=True, font_name=FONT_CN)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.72), Inches(10), Inches(0.35),
                     subtitle, font_size=13, color=RGBColor(0xDB, 0xEA, 0xFE))
    # 底部分隔线
    add_shape(slide, Inches(0), Inches(1.15), SLIDE_W, Pt(2), fill_color=ACCENT_LIGHT)


def add_footer(slide, page_num, total=15):
    """添加底部页码"""
    add_text_box(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.35),
                 f"{page_num} / {total}", font_size=10, color=TEXT_MUTED,
                 alignment=PP_ALIGN.RIGHT, font_name=FONT_EN)


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=TEXT_LIGHT,
                    bullet_char="▸", spacing=1.5):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_CN
        p.space_after = Pt(font_size * 0.6)
        p.line_spacing = Pt(font_size * spacing)

    return txBox


def add_table(slide, left, top, width, height, headers, rows, header_color=ACCENT_BLUE,
              col_widths=None):
    """添加表格"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # 设置列宽
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # 表头
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = TEXT_WHITE
            paragraph.font.name = FONT_CN
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 数据行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = TEXT_LIGHT
                paragraph.font.name = FONT_CN
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_SECONDARY if i % 2 == 0 else BG_CARD
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 表格边框
    for i in range(num_rows):
        for j in range(num_cols):
            cell = table.cell(i, j)
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(4)
            cell.margin_bottom = Pt(4)

    return table_shape


def add_stat_card(slide, left, top, width, height, number, label, accent=ACCENT_BLUE):
    """添加统计卡片"""
    card = add_rounded_shape(slide, left, top, width, height, fill_color=BG_CARD)
    # 顶部色条
    add_shape(slide, left + Inches(0.1), top, width - Inches(0.2), Pt(3), fill_color=accent)
    # 数字
    add_text_box(slide, left + Inches(0.1), top + Inches(0.2), width - Inches(0.2), Inches(0.6),
                 number, font_size=28, color=accent, bold=True, alignment=PP_ALIGN.CENTER,
                 font_name=FONT_EN)
    # 标签
    add_text_box(slide, left + Inches(0.1), top + Inches(0.8), width - Inches(0.2), Inches(0.4),
                 label, font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    return card


# ══════════════════════════════════════════════════════════════
#  Slide 1: 封面
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# 装饰线条
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, fill_color=ACCENT_BLUE)
add_shape(slide, Inches(0), Inches(3.4), SLIDE_W, Pt(1), fill_color=ACCENT_BLUE)
add_shape(slide, Inches(0), Inches(4.5), Inches(8), Pt(1), fill_color=BORDER_COLOR)

# 主标题
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(10), Inches(1.0),
             "Hermes Agent", font_size=52, color=ACCENT_BLUE, bold=True, font_name=FONT_EN)
add_text_box(slide, Inches(1.0), Inches(2.5), Inches(10), Inches(0.8),
             "源码深度分析", font_size=40, color=TEXT_WHITE, bold=True)

# 副标题信息
add_multiline_text(slide, Inches(1.0), Inches(3.6), Inches(10), Inches(1.0), [
    {"text": "Nous Research 开源自我改进 AI Agent", "size": 16, "color": TEXT_LIGHT},
    {"text": "仓库: github.com/NousResearch/hermes-agent  |  技术栈: Python 3.11+  |  许可证: MIT", "size": 12, "color": TEXT_MUTED},
])

# 右侧标签
tags = [
    ("闭环学习", ACCENT_GREEN),
    ("200+ 模型", ACCENT_BLUE),
    ("20+ 平台", ACCENT_AMBER),
    ("50+ 工具", ACCENT_ROSE),
]
for i, (tag, color) in enumerate(tags):
    y = Inches(4.8) + Inches(i * 0.45)
    shape = add_rounded_shape(slide, Inches(1.0), y, Inches(2.2), Inches(0.35), fill_color=color)
    add_text_box(slide, Inches(1.1), y + Inches(0.02), Inches(2.0), Inches(0.3),
                 tag, font_size=13, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 底部日期
add_text_box(slide, Inches(1.0), Inches(6.8), Inches(5), Inches(0.3),
             "分析日期: 2026-04-14  |  版本: main", font_size=11, color=TEXT_MUTED)
add_footer(slide, 1)


# ══════════════════════════════════════════════════════════════
#  Slide 2: 项目定位与核心差异
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "项目定位与核心差异", "Hermes Agent vs 其他 AI Agent 框架")
add_footer(slide, 2)

# 左侧: 核心差异
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5), Inches(0.4),
             "核心差异点", font_size=18, color=ACCENT_BLUE, bold=True)

features = [
    "闭环学习 — 从经验中自动创建 Skills，使用中改进，跨会话搜索记忆",
    "多模型支持 — OpenRouter 200+ 模型、Anthropic、OpenAI、GLM、Kimi、MiniMax",
    "多通道网关 — Telegram、Discord、飞书、微信、钉钉、企业微信、Slack 等 20+",
    "多终端后端 — 本地、Docker、SSH、Modal(Serverless)、Daytona、Singularity",
    "OpenClaw 兼容 — 内置 hermes claw migrate 迁移工具",
    "RL 训练内置 — Atropos 环境，支持 SWE-bench、TerminalBench 评估",
]
add_bullet_list(slide, Inches(0.6), Inches(1.9), Inches(5.8), Inches(4.5),
                features, font_size=13, spacing=1.6)

# 右侧: 架构流程
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5), Inches(0.4),
             "核心架构流程", font_size=18, color=ACCENT_CYAN, bold=True)

flow_items = [
    ("用户消息", ACCENT_BLUE),
    ("CLI / Gateway", BG_CARD),
    ("run_agent.py 主循环", BG_CARD),
    ("LLM API 调用", BG_CARD),
    ("工具执行 / Skill", BG_CARD),
    ("回复用户", ACCENT_GREEN),
]
for i, (item, color) in enumerate(flow_items):
    y = Inches(2.0) + Inches(i * 0.65)
    shape = add_rounded_shape(slide, Inches(7.2), y, Inches(4.5), Inches(0.5), fill_color=color)
    add_text_box(slide, Inches(7.3), y + Inches(0.05), Inches(4.3), Inches(0.4),
                 item, font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        # 箭头符号
        add_text_box(slide, Inches(9.2), y + Inches(0.45), Inches(0.5), Inches(0.25),
                     "▼", font_size=12, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# 侧边标注
add_text_box(slide, Inches(7.0), Inches(5.9), Inches(5.5), Inches(0.3),
             "记忆/上下文 ←→ Skill/MCP 双向支撑", font_size=11, color=TEXT_MUTED,
             alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  Slide 3: 代码规模统计
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "代码规模统计", "389K 行 Python · 15 个核心模块 · 50+ 内置工具")
add_footer(slide, 3)

# 统计卡片行
stats = [
    ("~389K", "Python 源码行数", ACCENT_BLUE),
    ("10,865", "run_agent.py 行数", ACCENT_AMBER),
    ("50+", "内置工具", ACCENT_GREEN),
    ("20+", "平台适配器", ACCENT_ROSE),
    ("26", "技能分类目录", ACCENT_CYAN),
    ("3,950", "飞书适配器行数", ACCENT_RED),
]
for i, (num, label, color) in enumerate(stats):
    col = i % 6
    x = Inches(0.5) + Inches(col * 2.1)
    y = Inches(1.5)
    add_stat_card(slide, x, y, Inches(1.95), Inches(1.2), num, label, color)

# 核心模块表格
add_text_box(slide, Inches(0.6), Inches(3.0), Inches(5), Inches(0.4),
             "核心模块代码量", font_size=16, color=ACCENT_BLUE, bold=True)

module_rows = [
    ("run_agent.py", "10,865", "Agent 主循环"),
    ("cli.py", "10,013", "CLI 入口 + TUI"),
    ("gateway/run.py", "~9,000", "消息网关"),
    ("hermes_cli/config.py", "3,306", "配置管理"),
    ("hermes_cli/auth.py", "3,270", "认证管理"),
    ("hermes_cli/setup.py", "3,138", "初始化向导"),
    ("hermes_cli/models.py", "1,933", "模型管理"),
    ("gateway/platforms/feishu.py", "3,950", "飞书适配器"),
]
add_table(slide, Inches(0.5), Inches(3.5), Inches(7.5), Inches(3.5),
          ["文件", "行数", "职责"], module_rows,
          col_widths=[Inches(3.5), Inches(1.2), Inches(2.8)])

# 右侧: 关键指标
add_text_box(slide, Inches(8.5), Inches(3.0), Inches(4.5), Inches(0.4),
             "关键指标", font_size=16, color=ACCENT_CYAN, bold=True)

indicators = [
    ("编程语言", "Python 3.11+"),
    ("包管理", "uv / pip"),
    ("框架", "OpenAI SDK + Fire CLI"),
    ("测试", "pytest"),
    ("最大单文件", "run_agent.py (10.8K 行)"),
    ("平台适配器", "20+ 个"),
    ("工具解析器", "11 种模型"),
    ("记忆插件", "8 种后端"),
]
for i, (key, val) in enumerate(indicators):
    y = Inches(3.5) + Inches(i * 0.42)
    add_text_box(slide, Inches(8.5), y, Inches(2.0), Inches(0.35),
                 key, font_size=12, color=TEXT_MUTED)
    add_text_box(slide, Inches(10.3), y, Inches(2.5), Inches(0.35),
                 val, font_size=12, color=TEXT_LIGHT)


# ══════════════════════════════════════════════════════════════
#  Slide 4: 核心架构概览
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "核心架构概览", "四层架构: 用户接口 → Agent 核心 → 工具执行 → 基础设施")
add_footer(slide, 4)

layers = [
    ("用户接口层", ["CLI (TUI)", "Gateway (多平台)", "Web (React)", "ACP Server"],
     ACCENT_BLUE, Inches(1.4)),
    ("Agent 核心层", ["run_agent.py / AIAgent", "prompt_builder (系统提示组装)",
                   "context_compressor (上下文压缩)", "memory_manager (记忆管理)",
                   "smart_model_routing (智能路由)"],
     ACCENT_GREEN, Inches(2.8)),
    ("工具执行层", ["Terminal (6后端)", "Browser (CamoFox)", "MCP (协议桥)",
                   "Delegate (子Agent)", "Skills (技能)", "Web (搜索)"],
     ACCENT_AMBER, Inches(4.6)),
    ("基础设施层", ["Config (配置)", "Cron (调度)", "Plugins (插件)",
                   "Registry (工具注册)", "ACP (协议适配)", "RL Env (训练)"],
     ACCENT_ROSE, Inches(6.0)),
]

for layer_name, items, color, y in layers:
    # 层标签
    label_shape = add_rounded_shape(slide, Inches(0.4), y, Inches(1.8), Inches(1.1), fill_color=color)
    add_text_box(slide, Inches(0.5), y + Inches(0.3), Inches(1.6), Inches(0.5),
                 layer_name, font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 层内容卡片
    for j, item in enumerate(items):
        x = Inches(2.5) + Inches(j * 2.15)
        card = add_rounded_shape(slide, x, y + Inches(0.1), Inches(2.0), Inches(0.9), fill_color=BG_CARD)
        # 顶部色条
        add_shape(slide, x, y + Inches(0.1), Inches(2.0), Pt(3), fill_color=color)
        add_text_box(slide, x + Inches(0.1), y + Inches(0.3), Inches(1.8), Inches(0.6),
                     item, font_size=11, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
#  Slide 5: 核心引擎 - AIAgent 主循环
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "核心引擎详解", "run_agent.py — AIAgent 主循环 (10,865 行)")
add_footer(slide, 5)

# 左列: 对话循环流程
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5), Inches(0.4),
             "run_conversation() 核心流程", font_size=16, color=ACCENT_BLUE, bold=True)

steps = [
    "1. 用户消息 → 追加到 messages",
    "2. 记忆预取 → memory_manager.prefetch_all()",
    "3. 组装系统提示 → _build_system_prompt()",
    "4. 上下文压缩检查 → context_engine.should_compress()",
    "5. 调用 LLM → call_llm (含 failover/重试)",
    "6. 工具调用 → registry.dispatch() 分发",
    "7. 循环直至无工具调用 → 最终回复",
    "8. 记忆同步 → memory_manager.sync_all()",
]
add_bullet_list(slide, Inches(0.6), Inches(1.9), Inches(5.5), Inches(4.0),
                steps, font_size=13, spacing=1.5, bullet_char="●")

# 右列: 关键机制
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.4),
             "关键机制", font_size=16, color=ACCENT_CYAN, bold=True)

mechanisms = [
    ("Failover & 重试", "classify_api_error() 分类 + jittered_backoff() 退避\n自动切换备用模型/端点"),
    ("智能模型路由", "根据消息复杂度关键词选择 cheap vs strong 模型\n可配置双模型策略，节省成本"),
    ("Prompt Caching", "Anthropic 原生缓存支持\n自动标记静态部分（系统提示、技能索引）"),
]

for i, (title, desc) in enumerate(mechanisms):
    y = Inches(2.0) + Inches(i * 1.4)
    card = add_rounded_shape(slide, Inches(7.0), y, Inches(5.5), Inches(1.2), fill_color=BG_CARD)
    add_shape(slide, Inches(7.0), y, Pt(4), Inches(1.2), fill_color=ACCENT_CYAN)
    add_text_box(slide, Inches(7.3), y + Inches(0.1), Inches(5.0), Inches(0.35),
                 title, font_size=14, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, Inches(7.3), y + Inches(0.45), Inches(5.0), Inches(0.7),
                 desc, font_size=11, color=TEXT_MUTED)

# 底部提示
add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
             "注: run_conversation() 约 2300 行内联逻辑，LLM 调用/工具分发/上下文压缩均未拆分为独立方法",
             font_size=10, color=TEXT_MUTED)


# ══════════════════════════════════════════════════════════════
#  Slide 6: 工具系统
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "工具系统", "tools/registry.py 注册中心 + 50+ 内置工具 + 6 种终端后端")
add_footer(slide, 6)

# 左侧: 工具注册
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(6), Inches(0.4),
             "Terminal Tool — 6 种执行后端", font_size=16, color=ACCENT_BLUE, bold=True)

backends = [
    ("local", "默认", "直接本机执行", ACCENT_GREEN),
    ("docker", "TERMINAL_ENV=docker", "容器隔离", ACCENT_BLUE),
    ("modal", "TERMINAL_ENV=modal", "Serverless 云沙箱", ACCENT_AMBER),
    ("ssh", "TERMINAL_ENV=ssh", "远程 SSH", ACCENT_CYAN),
    ("daytona", "TERMINAL_ENV=daytona", "Daytona 开发环境", ACCENT_ROSE),
    ("singularity", "TERMINAL_ENV=singularity", "HPC 环境", ACCENT_RED),
]
for i, (name, env, desc, color) in enumerate(backends):
    y = Inches(1.9) + Inches(i * 0.55)
    # 色条
    add_shape(slide, Inches(0.6), y, Pt(4), Inches(0.45), fill_color=color)
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(1.5), Inches(0.35),
                 name, font_size=13, color=color, bold=True, font_name=FONT_EN)
    add_text_box(slide, Inches(2.5), y + Inches(0.05), Inches(2.2), Inches(0.35),
                 env, font_size=10, color=TEXT_MUTED, font_name=FONT_EN)
    add_text_box(slide, Inches(4.8), y + Inches(0.05), Inches(2.0), Inches(0.35),
                 desc, font_size=11, color=TEXT_LIGHT)

# 右侧: 核心工具一览
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.4),
             "核心工具一览", font_size=16, color=ACCENT_AMBER, bold=True)

tool_rows = [
    ("Delegate Tool", "子 Agent 委托，并行执行，最大深度 2"),
    ("Browser Tool", "CamoFox 反指纹浏览器自动化"),
    ("MCP Tool", "Model Context Protocol 协议桥接"),
    ("Memory Tool", "读写 MEMORY.md / USER.md"),
    ("Skills Tool", "技能发现、创建、执行"),
    ("Session Search", "FTS5 全文搜索历史会话"),
    ("Cron Tools", "定时任务 CRUD 管理"),
    ("Vision Tool", "图像分析与识别"),
    ("TTS Tool", "语音合成"),
]
for i, (tool, desc) in enumerate(tool_rows):
    y = Inches(1.9) + Inches(i * 0.48)
    add_text_box(slide, Inches(7.0), y, Inches(2.3), Inches(0.4),
                 tool, font_size=12, color=ACCENT_AMBER, bold=True, font_name=FONT_EN)
    add_text_box(slide, Inches(9.3), y, Inches(3.5), Inches(0.4),
                 desc, font_size=11, color=TEXT_MUTED)

# 底部: 注册模式说明
add_rounded_shape(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.2), fill_color=BG_CARD)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(11.8), Inches(1.0),
             "工具注册模式: 每个 tools/*.py 在模块级调用 registry.register()，import 即注册 — 零配置发现\n"
             "导入链: tools/registry.py ← tools/*.py ← model_tools.py ← run_agent.py / cli.py / batch_runner.py",
             font_size=12, color=TEXT_MUTED)


# ══════════════════════════════════════════════════════════════
#  Slide 7: Skills 系统
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "Skills 系统", "Agent 的程序性记忆 — 捕获「如何做特定任务」的经验")
add_footer(slide, 7)

# 左侧: Skill 概念
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.4),
             "声明式记忆 vs 程序性记忆", font_size=16, color=ACCENT_BLUE, bold=True)

# 记忆对比卡片
card1 = add_rounded_shape(slide, Inches(0.6), Inches(1.9), Inches(2.6), Inches(2.0), fill_color=BG_CARD)
add_shape(slide, Inches(0.6), Inches(1.9), Inches(2.6), Pt(3), fill_color=ACCENT_BLUE)
add_text_box(slide, Inches(0.7), Inches(2.05), Inches(2.4), Inches(0.35),
             "声明式记忆 (是什么)", font_size=13, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.7), Inches(2.4), Inches(2.3), Inches(1.2),
                ["MEMORY.md", "USER.md", "宽泛的、事实性知识"],
                font_size=11, spacing=1.3, bullet_char="·")

card2 = add_rounded_shape(slide, Inches(3.5), Inches(1.9), Inches(2.6), Inches(2.0), fill_color=BG_CARD)
add_shape(slide, Inches(3.5), Inches(1.9), Inches(2.6), Pt(3), fill_color=ACCENT_GREEN)
add_text_box(slide, Inches(3.6), Inches(2.05), Inches(2.4), Inches(0.35),
             "程序性记忆 (怎么做)", font_size=13, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(3.6), Inches(2.4), Inches(2.3), Inches(1.2),
                ["Skills/", "SKILL.md", "具体的、可操作流程"],
                font_size=11, spacing=1.3, bullet_char="·")

# Skill 生命周期
add_text_box(slide, Inches(0.6), Inches(4.2), Inches(6), Inches(0.4),
             "Skill 生命周期", font_size=16, color=ACCENT_CYAN, bold=True)

lifecycle = [
    ("创建", "Agent 从复杂任务经验中自动创建\n或用户手动编写", ACCENT_BLUE),
    ("发现", "prompt_builder 扫描所有 SKILL.md\n构建索引 → 匹配用户消息", ACCENT_GREEN),
    ("执行", "作为系统提示的一部分\n指导 Agent 行为", ACCENT_AMBER),
    ("改进", "Agent 使用中通过\nskill_manager 编辑优化", ACCENT_ROSE),
    ("分享", "通过 Skills Hub\n发布到社区", ACCENT_CYAN),
]
for i, (phase, desc, color) in enumerate(lifecycle):
    x = Inches(0.6) + Inches(i * 2.4)
    card = add_rounded_shape(slide, x, Inches(4.7), Inches(2.2), Inches(1.5), fill_color=BG_CARD)
    add_shape(slide, x, Inches(4.7), Inches(2.2), Pt(3), fill_color=color)
    add_text_box(slide, x + Inches(0.1), Inches(4.85), Inches(2.0), Inches(0.35),
                 phase, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(5.2), Inches(2.0), Inches(0.9),
                 desc, font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# 右侧: Skills Hub + 安全
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.4),
             "Skills Hub 社区市场", font_size=16, color=ACCENT_AMBER, bold=True)

add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(2.0), [
    "兼容 agentskills.io 开放标准",
    "搜索/安装社区技能",
    "安全扫描 (skills_guard.py) — 检测恶意代码",
    "依赖检查 + 版本管理",
], font_size=13, spacing=1.5)

add_text_box(slide, Inches(7.0), Inches(3.5), Inches(5.5), Inches(0.4),
             "内置技能分类 (26 类)", font_size=16, color=ACCENT_ROSE, bold=True)

skill_cats = [
    "software-development · devops · data-science",
    "research · mlops · productivity · email",
    "mcp · red-teaming · gaming · smart-home",
    "social-media · creative · ...",
]
add_bullet_list(slide, Inches(7.0), Inches(4.0), Inches(5.5), Inches(2.0),
                skill_cats, font_size=12, spacing=1.4, bullet_char="·")


# ══════════════════════════════════════════════════════════════
#  Slide 8: Gateway 多平台适配
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "Gateway 多平台适配", "20+ 平台适配器 · 飞书 3,950 行 (最大平台文件)")
add_footer(slide, 8)

# 平台表格
platform_rows = [
    ("Telegram", "2,786", "语音转写、贴纸缓存、内联按钮"),
    ("Discord", "2,963", "Slash 命令、线程、嵌入消息"),
    ("Slack", "1,670", "Socket Mode、Block Kit"),
    ("飞书", "3,950", "事件订阅、卡片消息、群聊"),
    ("微信", "1,829", "桥接模式"),
    ("企业微信", "1,430", "WeCom 适配"),
    ("钉钉", "333", "DingTalk"),
    ("WhatsApp", "989", "桥接模式、语音备忘录"),
    ("Signal", "825", "signal-cli 桥接"),
    ("Matrix", "2,005", "Element 生态"),
    ("Email", "625", "SMTP/IMAP"),
    ("SMS", "373", "短信"),
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.0),
          ["平台", "代码行数", "特性"], platform_rows,
          col_widths=[Inches(1.5), Inches(1.2), Inches(4.8)])

# 右侧: 消息路由流程
add_text_box(slide, Inches(8.5), Inches(1.4), Inches(4.5), Inches(0.4),
             "消息路由流程", font_size=16, color=ACCENT_CYAN, bold=True)

route_steps = [
    "平台消息 → PlatformAdapter.on_message()",
    "GatewayRunner._handle_message()",
    "Hooks.pre_process() 前置钩子",
    "DM 配对检查 (pairing.py)",
    "Session 获取/创建",
    "AIAgent.run_conversation()",
    "Hooks.post_process() 后置钩子",
    "PlatformAdapter.send_message()",
]
for i, step in enumerate(route_steps):
    y = Inches(1.9) + Inches(i * 0.52)
    color = ACCENT_CYAN if i in [0, 5, 7] else TEXT_MUTED
    add_text_box(slide, Inches(8.5), y, Inches(4.3), Inches(0.45),
                 f"{i+1}. {step}", font_size=11, color=color)

# 安全提示
add_rounded_shape(slide, Inches(8.5), Inches(6.2), Inches(4.3), Inches(0.7), fill_color=BG_CARD)
add_text_box(slide, Inches(8.6), Inches(6.25), Inches(4.1), Inches(0.6),
             "DM 配对机制: 未配对用户 → 返回配对码 → CLI 确认 → 后续放行",
             font_size=11, color=ACCENT_AMBER)


# ══════════════════════════════════════════════════════════════
#  Slide 9: 记忆与上下文管理
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "记忆与上下文管理", "MemoryManager 编排 · 插件化后端 · 8 种外部记忆 · 上下文压缩")
add_footer(slide, 9)

# 左侧: 记忆管理器架构
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.4),
             "MemoryManager 架构", font_size=16, color=ACCENT_BLUE, bold=True)

mem_features = [
    "内置 Provider 永远存在，不可移除",
    "最多 1 个外部 Provider — 防止 Schema 膨胀",
    "故障隔离: 一个 Provider 失败不阻塞另一个",
    "记忆上下文隔离: <memory-context> 标签防止误读",
    "FTS5 全文搜索引擎索引所有历史会话",
]
add_bullet_list(slide, Inches(0.6), Inches(1.9), Inches(5.5), Inches(2.5),
                mem_features, font_size=12, spacing=1.5)

# 外部记忆插件
add_text_box(slide, Inches(0.6), Inches(3.9), Inches(5.5), Inches(0.4),
             "外部记忆插件", font_size=16, color=ACCENT_GREEN, bold=True)

mem_plugins = [
    ("Mem0", "智能记忆管理"),
    ("Supermemory", "云端记忆"),
    ("Holographic", "本地向量存储"),
    ("RetainDB", "数据库记忆"),
    ("Honcho", "方言式用户建模"),
    ("Hindsight", "回溯式记忆"),
    ("Byterover", "Byterover 集成"),
    ("OpenViking", "OpenViking"),
]
for i, (name, desc) in enumerate(mem_plugins):
    col = i % 4
    row = i // 4
    x = Inches(0.6) + Inches(col * 1.4)
    y = Inches(4.4) + Inches(row * 0.6)
    card = add_rounded_shape(slide, x, y, Inches(1.3), Inches(0.5), fill_color=BG_CARD)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.05), Inches(1.2), Inches(0.4),
                 name, font_size=10, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER,
                 font_name=FONT_EN)

# 右侧: 上下文压缩
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.4),
             "上下文压缩策略", font_size=16, color=ACCENT_AMBER, bold=True)

comp_flow = [
    "受保护头部 (3 条消息)",
    "待压缩中间部分",
    "   ↓ 1. 裁剪旧工具输出",
    "   ↓ 2. 辅助模型摘要",
    "   ↓ 3. 迭代更新摘要",
    "受保护尾部 (20 条消息)",
    "   ↓",
    "[头部] + [压缩摘要] + [尾部]",
]
for i, step in enumerate(comp_flow):
    y = Inches(1.9) + Inches(i * 0.4)
    color = ACCENT_AMBER if "↓" in step else TEXT_LIGHT
    add_text_box(slide, Inches(7.0), y, Inches(5.0), Inches(0.35),
                 step, font_size=11, color=color)

# 关键设计点
add_text_box(slide, Inches(7.0), Inches(5.2), Inches(5.5), Inches(0.4),
             "摘要质量保障", font_size=14, color=ACCENT_ROSE, bold=True)

quality_points = [
    "摘要前缀隔离: [CONTEXT COMPACTION — REFERENCE ONLY]",
    "失败冷却: 摘要失败后 600 秒内不重试",
    "Token 预算: summary_tokens = max(2000, min(budget, 12000))",
    "可替换引擎: config.yaml → context.engine 字段",
]
add_bullet_list(slide, Inches(7.0), Inches(5.6), Inches(5.5), Inches(1.5),
                quality_points, font_size=11, spacing=1.3, bullet_char="·")


# ══════════════════════════════════════════════════════════════
#  Slide 10: CLI 与配置体系
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "CLI 与配置体系", "cli.py (10K行) + hermes_cli/ (47个文件) + Smart Model Routing")
add_footer(slide, 10)

# 左侧: CLI 子命令
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(6), Inches(0.4),
             "CLI 子命令体系", font_size=16, color=ACCENT_BLUE, bold=True)

cli_rows = [
    ("hermes", "cli.py (TUI)", "启动交互对话"),
    ("hermes model", "models.py (1,933行)", "模型/Provider 选择"),
    ("hermes tools", "tools_config.py", "工具启用/禁用"),
    ("hermes config", "config.py (3,306行)", "配置管理"),
    ("hermes gateway", "gateway.py", "网关管理"),
    ("hermes setup", "setup.py (3,138行)", "初始化向导"),
    ("hermes auth", "auth.py (3,270行)", "认证管理"),
    ("hermes skills", "skills_config.py", "技能管理"),
    ("hermes cron", "cron.py", "定时任务"),
    ("hermes doctor", "doctor.py", "诊断问题"),
    ("hermes claw migrate", "claw.py", "OpenClaw 迁移"),
]
add_table(slide, Inches(0.5), Inches(1.9), Inches(6.5), Inches(4.5),
          ["命令", "模块", "功能"], cli_rows,
          col_widths=[Inches(2.2), Inches(2.0), Inches(2.3)])

# 右侧: 配置 + 模型路由
add_text_box(slide, Inches(7.5), Inches(1.4), Inches(5.3), Inches(0.4),
             "配置加载优先级", font_size=16, color=ACCENT_CYAN, bold=True)

priority = [
    "1. 命令行参数",
    "2. 环境变量 (HERMES_*)",
    "3. ~/.hermes/config.yaml",
    "4. ~/.hermes/.env",
    "5. 项目根目录 .env",
    "6. 内置默认值",
]
for i, item in enumerate(priority):
    y = Inches(1.9) + Inches(i * 0.4)
    color = ACCENT_CYAN if i < 3 else TEXT_MUTED
    add_text_box(slide, Inches(7.5), y, Inches(5.0), Inches(0.35),
                 item, font_size=12, color=color)

# Smart Model Routing
add_text_box(slide, Inches(7.5), Inches(4.4), Inches(5.3), Inches(0.4),
             "Smart Model Routing", font_size=16, color=ACCENT_AMBER, bold=True)

add_rounded_shape(slide, Inches(7.5), Inches(4.9), Inches(5.0), Inches(2.0), fill_color=BG_CARD)
add_multiline_text(slide, Inches(7.7), Inches(5.0), Inches(4.6), Inches(1.8), [
    {"text": "双模型策略", "size": 13, "color": ACCENT_AMBER, "bold": True},
    {"text": "简单问候/闲聊 → cheap 模型 (Haiku)", "size": 11, "color": TEXT_LIGHT},
    {"text": "复杂编码/调试 → strong 模型 (Sonnet)", "size": 11, "color": TEXT_LIGHT},
    {"text": "", "size": 6, "color": TEXT_MUTED},
    {"text": "关键词触发: debug, implement, refactor,", "size": 10, "color": TEXT_MUTED},
    {"text": "analyze, architecture, design, terminal...", "size": 10, "color": TEXT_MUTED},
])


# ══════════════════════════════════════════════════════════════
#  Slide 11: 插件与 ACP 适配
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "插件与 ACP 适配", "插件自动发现 · ACP 协议 · MCP 集成 · 11 种模型解析器")
add_footer(slide, 11)

# 左上: 插件系统
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.4),
             "插件系统", font_size=16, color=ACCENT_BLUE, bold=True)

plugin_steps = [
    "1. 扫描 plugins/ 目录",
    "2. 读取 plugin.yaml 元数据",
    "3. 配置中指定激活插件",
    "4. import 插件模块",
    "5. 调用 register() 注册",
    "6. 注入到对应管理器",
]
add_bullet_list(slide, Inches(0.6), Inches(1.9), Inches(5.5), Inches(2.5),
                plugin_steps, font_size=12, spacing=1.4, bullet_char="→")

add_text_box(slide, Inches(0.6), Inches(4.3), Inches(5.5), Inches(0.4),
             "插件类型", font_size=14, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(0.6), Inches(4.7), Inches(5.5), Inches(1.0), [
    "context_engine/ → 替换默认压缩引擎",
    "memory/ → 替换/增强记忆后端",
], font_size=12, spacing=1.4, bullet_char="●")

# 右上: ACP vs MCP
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.4),
             "ACP vs MCP 定位", font_size=16, color=ACCENT_AMBER, bold=True)

acp_mcp = [
    ("ACP (Agent Client Protocol)", [
        "Hermes 作为服务端，被其他系统调用",
        "支持: initialize / new_session / prompt / fork",
        "用途: IDE 插件、编排系统集成",
    ], ACCENT_BLUE),
    ("MCP (Model Context Protocol)", [
        "Hermes 作为客户端，调用外部工具",
        "传输: stdio (本地) / SSE (远程)",
        "支持: OAuth 2.0 认证",
    ], ACCENT_GREEN),
]
for i, (title, items, color) in enumerate(acp_mcp):
    y = Inches(1.9) + Inches(i * 2.0)
    card = add_rounded_shape(slide, Inches(7.0), y, Inches(5.5), Inches(1.8), fill_color=BG_CARD)
    add_shape(slide, Inches(7.0), y, Inches(5.5), Pt(3), fill_color=color)
    add_text_box(slide, Inches(7.2), y + Inches(0.1), Inches(5.1), Inches(0.35),
                 title, font_size=13, color=color, bold=True)
    add_bullet_list(slide, Inches(7.2), y + Inches(0.5), Inches(5.1), Inches(1.2),
                    items, font_size=11, spacing=1.3, bullet_char="·")

# 底部: 工具调用解析器
add_text_box(slide, Inches(0.6), Inches(5.8), Inches(12), Inches(0.4),
             "11 种模型工具调用解析器", font_size=14, color=ACCENT_ROSE, bold=True)

parsers = ["Hermes", "Llama", "Qwen", "Qwen3-Coder", "DeepSeek V3", "DeepSeek V3.1",
           "GLM-4.5", "GLM-4.7", "Kimi K2", "Mistral", "LongCat"]
for i, parser in enumerate(parsers):
    col = i % 6
    row = i // 6
    x = Inches(0.6) + Inches(col * 2.1)
    y = Inches(6.3) + Inches(row * 0.45)
    card = add_rounded_shape(slide, x, y, Inches(1.95), Inches(0.38), fill_color=BG_CARD)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.03), Inches(1.85), Inches(0.32),
                 parser, font_size=10, color=ACCENT_ROSE, alignment=PP_ALIGN.CENTER,
                 font_name=FONT_EN)


# ══════════════════════════════════════════════════════════════
#  Slide 12: 调度与自动化
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "调度与自动化", "Cron 定时调度 · 委托子 Agent · RL 训练环境 · 批量运行")
add_footer(slide, 12)

# 三列布局
# 列1: Cron
add_text_box(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(0.4),
             "Cron 调度", font_size=16, color=ACCENT_BLUE, bold=True)

cron_items = [
    "每 60 秒 tick 检查到期作业",
    "文件锁防止并发执行",
    "自然语言定义定时任务",
    "静默标记: [SILENT] 仅保存",
    "支持 17+ 投递平台",
    "jobs.json 持久化存储",
]
add_bullet_list(slide, Inches(0.5), Inches(1.9), Inches(3.8), Inches(3.0),
                cron_items, font_size=12, spacing=1.5, bullet_char="●")

# 列2: Delegate
add_text_box(slide, Inches(4.7), Inches(1.4), Inches(3.8), Inches(0.4),
             "委托子 Agent", font_size=16, color=ACCENT_GREEN, bold=True)

delegate_items = [
    "隔离上下文 (独立对话历史)",
    "独立终端会话",
    "受限工具集 (禁递归/记忆/消息)",
    "ThreadPoolExecutor 并行执行",
    "最大并发: 默认 3 (可配置)",
    "最大嵌套深度: 2",
]
add_bullet_list(slide, Inches(4.7), Inches(1.9), Inches(3.8), Inches(3.0),
                delegate_items, font_size=12, spacing=1.5, bullet_char="●")

# 列3: RL Training
add_text_box(slide, Inches(8.9), Inches(1.4), Inches(4.0), Inches(0.4),
             "RL 训练环境", font_size=16, color=ACCENT_AMBER, bold=True)

rl_items = [
    "Atropos RL 训练框架集成",
    "SWE-bench 自动化评估",
    "TerminalBench v2 基准测试",
    "批量轨迹生成 (batch_runner.py)",
    "轨迹压缩 (trajectory_compressor.py)",
    "可复用 Agent Loop (agent_loop.py)",
]
add_bullet_list(slide, Inches(8.9), Inches(1.9), Inches(4.0), Inches(3.0),
                rl_items, font_size=12, spacing=1.5, bullet_char="●")

# 底部: 并发控制说明
add_rounded_shape(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.3), fill_color=BG_CARD)
add_multiline_text(slide, Inches(0.7), Inches(5.6), Inches(11.8), Inches(1.1), [
    {"text": "安全边界 — DELEGATE_BLOCKED_TOOLS", "size": 14, "color": ACCENT_RED, "bold": True},
    {"text": "delegate_task (禁止递归) · clarify (禁止用户交互) · memory (禁止写共享记忆) · send_message (禁止跨平台) · execute_code (应逐步推理)", "size": 12, "color": TEXT_MUTED},
    {"text": "全局工具线程池: 默认 128 workers，运行时可动态调整", "size": 12, "color": TEXT_MUTED},
])


# ══════════════════════════════════════════════════════════════
#  Slide 13: 安全策略纵深
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "安全策略纵深", "输入防护 · 执行防护 · 输出防护 — 三层安全防线")
add_footer(slide, 13)

# 三列: 输入/执行/输出
security_layers = [
    ("输入防护", ACCENT_BLUE, [
        ("Prompt 注入扫描", "10+ 威胁模式匹配\n不可见 Unicode 字符检测"),
        ("上下文文件扫描", "AGENTS.md / SOUL.md 加载前扫描\n替换为警告文本"),
        ("URL 安全检查", "恶意域名黑名单\n私有 IP 检测 (SSRF 防护)"),
    ]),
    ("执行防护", ACCENT_GREEN, [
        ("命令审批系统", "allow / ask / deny 三种模式\n10+ 危险命令模式匹配"),
        ("沙箱隔离", "6 种后端 (Docker/Modal/SSH/...)\n容器默认无网络访问"),
        ("工具限制", "子 Agent 受限工具集\n最大嵌套深度 2"),
    ]),
    ("输出防护", ACCENT_AMBER, [
        ("路径安全", "realpath 校验\n工作目录限制 + 二次确认"),
        ("凭证保护", "API Key 加密存储\nKey 池轮转避免限流"),
        ("技能安全扫描", "安装前扫描恶意代码\n网络外传 / 文件破坏检测"),
    ]),
]

for col_idx, (layer_name, color, items) in enumerate(security_layers):
    x = Inches(0.4) + Inches(col_idx * 4.3)

    # 层标题
    header = add_rounded_shape(slide, x, Inches(1.4), Inches(4.0), Inches(0.5), fill_color=color)
    add_text_box(slide, x + Inches(0.1), Inches(1.45), Inches(3.8), Inches(0.4),
                 layer_name, font_size=16, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 安全项
    for i, (title, desc) in enumerate(items):
        y = Inches(2.1) + Inches(i * 1.5)
        card = add_rounded_shape(slide, x, y, Inches(4.0), Inches(1.3), fill_color=BG_CARD)
        add_shape(slide, x, y, Pt(4), Inches(1.3), fill_color=color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.35),
                     title, font_size=13, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.45), Inches(3.6), Inches(0.8),
                     desc, font_size=10, color=TEXT_MUTED)

# 底部: DM 配对安全
add_rounded_shape(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6), fill_color=BG_CARD)
add_text_box(slide, Inches(0.7), Inches(6.65), Inches(11.8), Inches(0.5),
             "DM 配对安全: 用户发消息 → 检查配对状态 → 未配对返回配对码 → CLI 确认 → 后续放行  |  Tirith 安全框架集成",
             font_size=11, color=ACCENT_RED)


# ══════════════════════════════════════════════════════════════
#  Slide 14: 与 OpenClaw 对比
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_header_bar(slide, "与 OpenClaw 架构对比", "Hermes Agent (Python) vs OpenClaw (TypeScript)")
add_footer(slide, 14)

compare_rows = [
    ("语言", "TypeScript / Node.js", "Python"),
    ("定位", "多通道 AI 网关", "自我改进 AI Agent"),
    ("核心循环", "Gateway → Agent → Channel", "CLI/Gateway → AIAgent → LLM"),
    ("工具系统", "内置 + MCP", "内置 Registry + MCP + Skills"),
    ("记忆", "MEMORY.md + 外部 SDK", "MemoryManager + 8 种插件后端 + FTS5"),
    ("平台数", "30+", "20+"),
    ("模型", "20+ Provider", "OpenRouter 200+ / 直连"),
    ("学习闭环", "无", "Skill 自动创建 + 自改进 + 搜索"),
    ("RL 训练", "无", "内置 Atropos 环境"),
    ("子 Agent", "sessions_spawn", "delegate_tool (并行)"),
    ("MCP", "✅", "✅"),
    ("飞书", "一等公民", "平台适配器 (3,950 行)"),
    ("代码量", "~613K 行 TS", "~389K 行 Python"),
]
add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
          ["维度", "OpenClaw", "Hermes Agent"], compare_rows,
          col_widths=[Inches(2.0), Inches(4.5), Inches(5.8)])


# ══════════════════════════════════════════════════════════════
#  Slide 15: 关键设计哲学 / 结束
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# 装饰线条
add_shape(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, fill_color=ACCENT_BLUE)

add_text_box(slide, Inches(1.0), Inches(0.8), Inches(10), Inches(0.6),
             "关键设计哲学", font_size=30, color=ACCENT_BLUE, bold=True)

philosophies = [
    ("单文件巨石", "run_agent.py 10K+ 行承载完整 Agent 主循环 — Python 风格的「一个文件搞定」", ACCENT_BLUE),
    ("注册模式", "工具通过 registry.register() 在 import 时自注册，零配置发现", ACCENT_GREEN),
    ("插件化扩展", "记忆、上下文引擎均可通过插件替换，放文件即生效", ACCENT_AMBER),
    ("Skill 即程序记忆", "区别于声明式记忆 (MEMORY.md)，Skill 是过程性知识", ACCENT_ROSE),
    ("Serverless 友好", "Modal/Daytona 后端支持空闲休眠、按需唤醒", ACCENT_CYAN),
    ("安全纵深", "Prompt 注入扫描 → 命令审批 → 沙箱隔离 → 路径安全检查", ACCENT_RED),
]

for i, (title, desc, color) in enumerate(philosophies):
    y = Inches(1.6) + Inches(i * 0.85)
    # 序号圆
    add_rounded_shape(slide, Inches(1.0), y, Inches(0.45), Inches(0.45), fill_color=color)
    add_text_box(slide, Inches(1.05), y + Inches(0.03), Inches(0.4), Inches(0.4),
                 str(i + 1), font_size=16, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 font_name=FONT_EN)
    # 内容
    add_text_box(slide, Inches(1.6), y, Inches(2.0), Inches(0.4),
                 title, font_size=16, color=color, bold=True)
    add_text_box(slide, Inches(1.6), y + Inches(0.35), Inches(10), Inches(0.4),
                 desc, font_size=12, color=TEXT_MUTED)

# 分隔线
add_shape(slide, Inches(1.0), Inches(6.8), Inches(8), Pt(1), fill_color=BORDER_COLOR)

# 结束语
add_text_box(slide, Inches(1.0), Inches(6.9), Inches(10), Inches(0.4),
             "Hermes Agent — 让 AI Agent 从经验中学习，自我改进  |  github.com/NousResearch/hermes-agent",
             font_size=12, color=TEXT_MUTED)
add_footer(slide, 15)


# ── 保存文件 ──
output_path = os.path.join(os.path.dirname(__file__), "..", "..",
                           "Hermes-Agent源码深度分析.pptx")
output_path = os.path.abspath(output_path)
prs.save(output_path)
print(f"演示文稿已生成: {output_path}")
print(f"共 {len(prs.slides)} 页")
