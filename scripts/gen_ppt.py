# -*- coding: utf-8 -*-
"""
文件名：gen_ppt.py
描述：基于 Hermes-vs-OpenClaw 对比分析文档生成 PowerPoint 演示文稿
作者：Claude Code
创建日期：2026-04-14
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色主题 ──
PRIMARY = RGBColor(0x1A, 0x56, 0xDB)       # 主蓝
PRIMARY_DARK = RGBColor(0x0F, 0x3D, 0xA8)   # 深蓝
ACCENT_HERMES = RGBColor(0xE8, 0x6B, 0x2C)  # Hermes 橙
ACCENT_OPENCLAW = RGBColor(0x0E, 0xA5, 0x6B)  # OpenClaw 绿
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)       # 浅灰背景
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)       # 深色文字
TEXT_MED = RGBColor(0x4A, 0x55, 0x68)        # 中灰文字
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)      # 白色文字
HEADER_BG = RGBColor(0x1A, 0x56, 0xDB)      # 表头蓝
ROW_ALT = RGBColor(0xEE, 0xF2, 0xFF)        # 交替行浅蓝
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER_COLOR = RGBColor(0xD1, 0xD5, 0xDB)


# ── 工具函数 ──

def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=14, bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """设置形状文字"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_text_paragraph(tf, text, font_size=14, bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT, space_before=Pt(4), font_name="Microsoft YaHei"):
    """向文本框追加段落"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bullet_points(tf, items, font_size=13, color=TEXT_DARK, bullet_char="•", indent_level=0, font_name="Microsoft YaHei"):
    """添加项目符号列表"""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.level = indent_level
        p.space_before = Pt(4)


def add_comparison_table(slide, left, top, width, height, headers, rows):
    """
    添加对比表格
    headers: ["维度", "Hermes Agent", "OpenClaw"]
    rows: list of lists
    """
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # 设置列宽比例
    col_widths = [Inches(1.8), Inches(3.1), Inches(3.1)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 数据行
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            bg = ROW_ALT if r_idx % 2 == 0 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_DARK
                p.font.name = "Microsoft YaHei"
                if c_idx == 0:
                    p.font.bold = True
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_title_subtitle(slide, title_text, subtitle_text=None):
    """添加标题和副标题装饰条"""
    # 顶部蓝色装饰条
    add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), PRIMARY)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(8.6), Inches(0.6))
    set_text(title_box, title_text, font_size=24, bold=True, color=PRIMARY)

    # 副标题
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.85), Inches(8.6), Inches(0.4))
        set_text(sub_box, subtitle_text, font_size=13, color=TEXT_MED)

    # 底部分隔线
    add_shape(slide, Inches(0.7), Inches(1.25), Inches(8.6), Inches(0.02), PRIMARY)


def add_key_insight(slide, text, top=Inches(6.5)):
    """添加底部关键洞察条"""
    box = add_rounded_rect(slide, Inches(0.7), top, Inches(8.6), Inches(0.6), RGBColor(0xEF, 0xF6, 0xFF))
    tf = set_text(box, f"💡 {text}", font_size=11, color=PRIMARY_DARK)
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT


# ── 幻灯片内容 ──

def slide_cover(prs):
    """封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, WHITE)

    # 左侧大色块
    add_shape(slide, Inches(0), Inches(0), Inches(4.5), Inches(7.5), PRIMARY)

    # 左侧标题
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(3.5), Inches(2))
    tf = set_text(box, "架构对比分析", font_size=36, bold=True, color=WHITE)
    add_text_paragraph(tf, "Architecture Comparison", font_size=16, color=RGBColor(0xBB, 0xD5, 0xFF))

    # 左侧项目标签
    hermes_tag = add_rounded_rect(slide, Inches(0.5), Inches(3.5), Inches(1.6), Inches(0.5), ACCENT_HERMES)
    set_text(hermes_tag, "Hermes Agent", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    vs_tag = add_rounded_rect(slide, Inches(2.3), Inches(3.5), Inches(0.5), Inches(0.5), WHITE)
    set_text(vs_tag, "VS", font_size=12, bold=True, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    openclaw_tag = add_rounded_rect(slide, Inches(3.0), Inches(3.5), Inches(1.3), Inches(0.5), ACCENT_OPENCLAW)
    set_text(openclaw_tag, "OpenClaw", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 右侧信息
    info_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.3), Inches(3.5))
    tf = set_text(info_box, "自我进化的 AI Agent 平台", font_size=16, bold=True, color=ACCENT_HERMES)
    add_text_paragraph(tf, "vs", font_size=14, color=TEXT_MED, alignment=PP_ALIGN.CENTER)
    add_text_paragraph(tf, "多渠道 AI 网关 / Agent 平台", font_size=16, bold=True, color=ACCENT_OPENCLAW)
    add_text_paragraph(tf, "", font_size=8)
    add_text_paragraph(tf, "对比版本", font_size=12, bold=True, color=TEXT_DARK)
    add_text_paragraph(tf, "Hermes Agent v0.9.0 / OpenClaw v2026.4.11", font_size=11, color=TEXT_MED)
    add_text_paragraph(tf, "", font_size=6)
    add_text_paragraph(tf, "分析日期：2026-04-14", font_size=11, color=TEXT_MED)

    # 底部一句话
    bottom = add_rounded_rect(slide, Inches(0.5), Inches(5.5), Inches(3.5), Inches(0.8), RGBColor(0x14, 0x47, 0xB1))
    set_text(bottom, '"全能 Agent"  vs  "万能网关"', font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_toc(prs):
    """目录"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "目录", "分析维度总览")

    toc_items = [
        ("01", "项目定位与核心理念"),
        ("02", "技术栈与工程规模"),
        ("03", "整体架构对比"),
        ("04", "Agent 核心引擎"),
        ("05", "消息渠道与多平台适配"),
        ("06", "工具系统与技能扩展"),
        ("07", "内存与上下文管理"),
        ("08", "安全策略"),
        ("09", "插件与扩展生态"),
        ("10", "部署与运维"),
        ("11", "AI 模型接入"),
        ("12", "多 Agent 配置与多通道互通"),
        ("13", "总结与适用场景"),
    ]

    # 两列布局
    for i, (num, title) in enumerate(toc_items):
        col = 0 if i < 7 else 1
        row = i if i < 7 else i - 7
        x = Inches(0.7 + col * 4.5)
        y = Inches(1.6 + row * 0.7)

        num_box = add_rounded_rect(slide, x, y, Inches(0.5), Inches(0.45), PRIMARY)
        set_text(num_box, num, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

        title_box = slide.shapes.add_textbox(x + Inches(0.6), y, Inches(3.5), Inches(0.45))
        set_text(title_box, title, font_size=13, color=TEXT_DARK)


def slide_positioning(prs):
    """项目定位与核心理念"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "01 项目定位与核心理念", "定位决定架构走向")

    # Hermes 卡片
    h_card = add_rounded_rect(slide, Inches(0.7), Inches(1.5), Inches(4.2), Inches(3.5), RGBColor(0xFF, 0xF5, 0xEE))
    tf = set_text(h_card, "Hermes Agent", font_size=18, bold=True, color=ACCENT_HERMES)
    h_card.text_frame.word_wrap = True
    add_text_paragraph(tf, "自我进化的 AI Agent 平台", font_size=14, bold=True, color=TEXT_DARK)
    add_text_paragraph(tf, "", font_size=6)
    add_text_paragraph(tf, '• 核心理念：闭环学习', font_size=12, color=TEXT_DARK)
    add_text_paragraph(tf, '  Agent 从经验中创建技能', font_size=11, color=TEXT_MED)
    add_text_paragraph(tf, '  使用中优化技能', font_size=11, color=TEXT_MED)
    add_text_paragraph(tf, '  搜索历史对话', font_size=11, color=TEXT_MED)
    add_text_paragraph(tf, "", font_size=4)
    add_text_paragraph(tf, '• 技术栈：Python 3.11+', font_size=12, color=TEXT_DARK)
    add_text_paragraph(tf, '• 许可证：MIT', font_size=12, color=TEXT_DARK)
    add_text_paragraph(tf, '• "一个会学习的 Agent"', font_size=12, bold=True, color=ACCENT_HERMES)

    # OpenClaw 卡片
    o_card = add_rounded_rect(slide, Inches(5.1), Inches(1.5), Inches(4.2), Inches(3.5), RGBColor(0xEE, 0xFB, 0xF5))
    tf = set_text(o_card, "OpenClaw", font_size=18, bold=True, color=ACCENT_OPENCLAW)
    o_card.text_frame.word_wrap = True
    add_text_paragraph(tf, "多渠道 AI 网关 / Agent 平台", font_size=14, bold=True, color=TEXT_DARK)
    add_text_paragraph(tf, "", font_size=6)
    add_text_paragraph(tf, '• 核心理念：网关优先', font_size=12, color=TEXT_DARK)
    add_text_paragraph(tf, '  统一连接消息平台与 AI 模型', font_size=11, color=TEXT_MED)
    add_text_paragraph(tf, '  Plugin-first 架构', font_size=11, color=TEXT_MED)
    add_text_paragraph(tf, "", font_size=4)
    add_text_paragraph(tf, '• 技术栈：TypeScript (ESM)', font_size=12, color=TEXT_DARK)
    add_text_paragraph(tf, '• 运行时：Node.js >= 22.14', font_size=12, color=TEXT_DARK)
    add_text_paragraph(tf, '• "一个连接一切的 AI 网关"', font_size=12, bold=True, color=ACCENT_OPENCLAW)

    add_key_insight(slide, "核心差异：Hermes 差异化在于自我学习能力，OpenClaw 差异化在于连接广度（30+ 消息平台、20+ AI 模型提供商）", top=Inches(5.4))


def slide_tech_stack(prs):
    """技术栈与工程规模"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "02 技术栈与工程规模", "Python vs TypeScript，两种工程哲学")

    headers = ["维度", "Hermes Agent", "OpenClaw"]
    rows = [
        ["主语言", "Python 3.11+", "TypeScript (ESM, strict)"],
        ["运行时", "CPython", "Node.js >= 22.14"],
        ["包管理", "pip / uv", "pnpm 10.32 (monorepo)"],
        ["测试框架", "Pytest (~3000 测试)", "Vitest 4.1.4"],
        ["核心源码", "~389K 行 Python", "~657K 行 TypeScript"],
        ["文件数量", "~200+ Python 文件", "~3800 TS 文件"],
        ["最大单文件", "run_agent.py: 10,865 行", "config/schema.ts: 27,117 行"],
        ["仓库类型", "单仓库", "Monorepo (pnpm workspace)"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4.8), headers, rows)

    add_key_insight(slide, "Hermes 以 Python 简洁换取更高代码密度，OpenClaw 以 TypeScript 类型安全换取更模块化组织")


def slide_architecture(prs):
    """整体架构对比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "03 整体架构对比", "大内核 vs 小内核，两种架构路径")

    # Hermes 架构
    h_label = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.45), ACCENT_HERMES)
    set_text(h_label, "Hermes Agent — 大内核 + 薄插件", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    h_layers = [
        ("用户界面层", "CLI / Gateway / Web / ACP"),
        ("Agent 核心", "run_agent.py (AIAgent 单文件)"),
        ("工具执行", "50+ 工具 (导入自注册) + MCP"),
        ("基础设施", "Config / Cron / Plugins / RL Training"),
    ]
    for i, (name, desc) in enumerate(h_layers):
        y = Inches(2.1 + i * 0.85)
        box = add_rounded_rect(slide, Inches(0.5), y, Inches(4.3), Inches(0.75), RGBColor(0xFF, 0xF5, 0xEE))
        tf = set_text(box, name, font_size=11, bold=True, color=ACCENT_HERMES)
        add_text_paragraph(tf, desc, font_size=10, color=TEXT_MED)

    # OpenClaw 架构
    o_label = add_rounded_rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.45), ACCENT_OPENCLAW)
    set_text(o_label, "OpenClaw — 小内核 + 厚插件", font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    o_layers = [
        ("渠道插件层", "111 扩展目录，25+ 消息平台"),
        ("Gateway 控制平面", "WebSocket/HTTP，60+ RPC 方法"),
        ("Agent 引擎", "pi-embedded-runner (内嵌 LLM)"),
        ("扩展系统", "Extensions / Plugin SDK / 53 内置"),
        ("基础设施", "Config (Zod) / Sessions / Daemon"),
    ]
    for i, (name, desc) in enumerate(o_layers):
        y = Inches(2.1 + i * 0.85)
        box = add_rounded_rect(slide, Inches(5.2), y, Inches(4.3), Inches(0.75), RGBColor(0xEE, 0xFB, 0xF5))
        tf = set_text(box, name, font_size=11, bold=True, color=ACCENT_OPENCLAW)
        add_text_paragraph(tf, desc, font_size=10, color=TEXT_MED)

    # 数据流差异
    flow_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(0.5))
    tf = set_text(flow_box, "关键差异：", font_size=11, bold=True, color=PRIMARY)
    add_text_paragraph(tf, "Hermes 同步循环（简单直接）vs OpenClaw 异步流（高并发天然优势）", font_size=11, color=TEXT_DARK)


def slide_agent_engine(prs):
    """Agent 核心引擎"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "04 Agent 核心引擎", "Agent 循环的设计哲学差异")

    headers = ["维度", "Hermes Agent", "OpenClaw"]
    rows = [
        ["核心文件", "run_agent.py (10,865 行)", "run.ts + attempt.ts (4,200+ 行)"],
        ["Agent 循环", "同步 while 循环 (~2300 行)", "异步流处理，分阶段组织"],
        ["上下文压缩", "保护前3 + 后20 + 中间摘要", "Session 管理 + 压缩支持"],
        ["智能路由", "关键词复杂度 → 廉/强模型", "模型回退链 (fallback)"],
        ["System Prompt", "9 阶段管道组装", "Agent 配置系统提示"],
        ["子 Agent", "Delegate: 最大深度2, 并发3", "sub-agent: 异步, 并发5"],
        ["流式输出", "支持 (streaming tool)", "支持 (WebSocket 流式)"],
        ["工具解析器", "11 个模型专用解析器", "标准 OpenAI 格式"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4.8), headers, rows)

    add_key_insight(slide, "Hermes 偏'大泥球'（集中但可调试），OpenClaw 偏分层模块化（清晰但需跨文件追踪）")


def slide_channels(prs):
    """消息渠道与多平台适配"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "05 消息渠道与多平台适配", "中国平台 vs 国际平台的覆盖差异")

    headers = ["维度", "Hermes Agent", "OpenClaw"]
    rows = [
        ["平台数量", "17+ 适配器", "25+ 渠道插件 (111 扩展)"],
        ["中国平台", "飞书、企微、微信、钉钉", "飞书、钉钉"],
        ["国际 IM", "TG/Discord/Slack/WA/Signal...", "TG/Discord/Slack/WA/Signal/IRC/Line..."],
        ["特殊渠道", "Email/SMS/HomeAssistant", "MS Teams/Web API/Companion"],
        ["移动端", "无原生 App", "Android / iOS / macOS"],
        ["适配器架构", "BasePlatformAdapter (82K)", "Channel Plugin SDK (100+ 导出)"],
        ["加载方式", "asyncio.gather 并行启动", "Plugin Loader 发现加载"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4.2), headers, rows)

    add_key_insight(slide, "Hermes 中国平台覆盖更强（微信、企微、钉钉、飞书），OpenClaw 国际平台广度领先 + 有原生移动端")


def slide_tools(prs):
    """工具系统与技能扩展"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "06 工具系统与技能扩展", "工具注册、终端执行、技能创建的差异")

    headers = ["维度", "Hermes Agent", "OpenClaw"]
    rows = [
        ["注册方式", "导入时自注册", "Plugin SDK 声明式"],
        ["工具数量", "50+ 工具文件", "内置 + 扩展 + MCP"],
        ["终端执行", "6 种后端 (含 Docker/SSH/Modal)", "内置 exec + Docker sandbox"],
        ["浏览器", "CamoFox 反指纹浏览器", "Playwright 自动化"],
        ["MCP 角色", "客户端 (连接外部)", "服务器 (暴露自身)"],
        ["技能格式", "SKILL.md (YAML+MD)", "SKILL.md (description/when)"],
        ["内置技能", "27 个类别", "53 个内置技能"],
        ["自动创建", "✅ Agent 自动创建技能", "❌ 无自动创建"],
        ["技能市场", "Skills Hub + 安全扫描", "无独立市场"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(5.0), headers, rows)

    add_key_insight(slide, "Hermes 技能系统是闭环的（使用+创建+优化），OpenClaw 更偏声明式触发", top=Inches(6.7))


def slide_memory(prs):
    """内存与上下文管理"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "07 内存与上下文管理", "记忆系统的深度与广度")

    headers = ["维度", "Hermes Agent", "OpenClaw"]
    rows = [
        ["内置记忆", "MEMORY.md + USER.md", "Memory Host SDK (向量+SQLite)"],
        ["搜索能力", "FTS5 全文搜索历史会话", "SQLite-vec 向量搜索 + LanceDB"],
        ["外部后端", "8 个 (Mem0/Supermemory...)", "无独立外部后端"],
        ["上下文压缩", "75% token 阈值触发", "Session 上下文压缩"],
        ["上下文引擎", "可插拔 (LCM DAG 等)", "内置 Context Engine"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(3.2), headers, rows)

    add_key_insight(slide, "Hermes 在记忆系统投入更深（8 种外部后端、FTS5 全文搜索），OpenClaw 在向量检索场景可能更优", top=Inches(5.0))


def slide_security(prs):
    """安全策略"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "08 安全策略", "纵深防御的不同侧重")

    headers = ["安全层", "Hermes Agent", "OpenClaw"]
    rows = [
        ["输入防护", "10+ 正则检测注入 + Unicode 检测", "上下文可见性控制"],
        ["执行防护", "命令审批 + 6 种沙箱后端", "4 层策略管道 + 沙箱隔离"],
        ["沙箱后端", "Docker/Modal/Daytona/SSH/Singularity/local", "Docker / SSH"],
        ["凭证管理", "加密存储 + 多密钥轮换", "加密 + 系统密钥链集成"],
        ["技能安全", "安装前扫描 (恶意/外泄/危险)", "Skill Scanner 静态分析"],
        ["HTTP 拒绝列表", "无", "默认阻止 exec/spawn/shell 等"],
        ["URL 安全", "恶意域名 + 私有 IP (SSRF)", "无特殊处理"],
        ["安全审计", "无内置命令", "openclaw doctor security"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4.8), headers, rows)

    add_key_insight(slide, "平手：Hermes 重输入侧防护和执行隔离，OpenClaw 重工具权限控制和运维安全")


def slide_plugins(prs):
    """插件与扩展生态"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "09 插件与扩展生态", "松散灵活 vs 标准化成熟")

    headers = ["维度", "Hermes Agent", "OpenClaw"]
    rows = [
        ["插件类型", "memory/context_engine", "Channel/Provider/Capability/Tool"],
        ["插件发现", "目录扫描 plugins/<type>/<name>", "Plugin Loader (2,212 行)"],
        ["插件 SDK", "抽象基类 (ABC)", "@openclaw/plugin-sdk (100+ 导出)"],
        ["扩展数量", "~10 个插件", "111 个扩展目录"],
        ["技能数量", "27 类内置 + 大量可选", "53 个内置"],
        ["插件市场", "Skills Hub (agentskills.io)", "无"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(3.6), headers, rows)

    add_key_insight(slide, "OpenClaw 插件生态更成熟标准化（Plugin SDK + 完整生命周期），Hermes 更灵活但松散")


def slide_deploy(prs):
    """部署与运维"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "10 部署与运维", "个人开发者场景 vs 生产级运维")

    headers = ["维度", "Hermes Agent", "OpenClaw"]
    rows = [
        ["部署模式", "CLI / Docker / ACP", "CLI / Daemon / Docker"],
        ["Docker Compose", "无官方 compose", "有 (Gateway + CLI, 健康检查)"],
        ["守护进程", "无", "systemd / launchd / schtasks"],
        ["健康检查", "无内置", "/healthz, /readyz, /metrics"],
        ["监控", "无内置", "Prometheus metrics 端点"],
        ["配套 App", "无", "Android / iOS / macOS"],
        ["配置验证", "Pydantic + YAML", "Zod 4.x + JSON Schema"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4.2), headers, rows)

    add_key_insight(slide, "OpenClaw 运维明显更成熟（Daemon + Prometheus + Docker Compose + 多平台 App），Hermes 偏研究场景")


def slide_ai_models(prs):
    """AI 模型接入"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "11 AI 模型接入", "广度与深度的不同策略")

    headers = ["维度", "Hermes Agent", "OpenClaw"]
    rows = [
        ["提供商数量", "10+ (含 OpenRouter 200+)", "~30 个 provider 插件"],
        ["模型路由", "智能路由 (复杂度检测)", "模型回退链 (fallback)"],
        ["认证方式", "API Key + OAuth + 多密钥轮换", "API Key + OAuth + Auth Profile"],
        ["中国模型", "GLM / Kimi / MiniMax / MiMo", "DeepSeek"],
        ["本地模型", "vLLM / SGLang", "Ollama"],
        ["特殊能力", "11 种模型专用解析器", "标准 OpenAI 格式"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(3.6), headers, rows)

    add_key_insight(slide, "Hermes 模型接入更广（OpenRouter 200+ 模型）且中国模型支持更深，OpenClaw provider 插件更多但各自维护")


def slide_multi_agent(prs):
    """多 Agent 配置与多通道互通"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "12 多 Agent 配置与多通道互通", "架构差异最大的维度")

    headers = ["维度", "Hermes Agent", "OpenClaw"]
    rows = [
        ["配置模式", "多 Profile (多进程)", "单进程多 Agent (配置数组)"],
        ["实例隔离", "独立进程/数据库/凭证", "共享 Gateway，独立工作空间"],
        ["路由粒度", "平台级", "9 层匹配 (用户/群组/角色级)"],
        ["跨 Agent 委派", "仅自身实例", "可指定不同 AgentId"],
        ["运行时控制", "无主动控制", "steer / kill / list"],
        ["跨通道 Session", "不支持", "identityLinks 统一"],
        ["子 Agent 并发", "3 (ThreadPool, 受 GIL)", "5 (异步并发)"],
        ["混合专家推理", "Mixture-of-Agents ✅", "无 ❌"],
    ]
    add_comparison_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(4.8), headers, rows)

    add_key_insight(slide, "多 Agent 多通道维度上 OpenClaw 明显领先，Hermes 优势在实例隔离彻底 + Mixture-of-Agents")


def slide_matrix(prs):
    """核心差异矩阵"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "13 核心差异矩阵", "各维度能力评分一览")

    # 星级评分表
    headers = ["能力维度", "Hermes", "OpenClaw", "优势方"]
    rows = [
        ["自我学习与进化", "★★★★★", "★★☆☆☆", "Hermes"],
        ["中国消息平台", "★★★★☆", "★★☆☆☆", "Hermes"],
        ["国际消息平台", "★★★☆☆", "★★★★★", "OpenClaw"],
        ["工具系统深度", "★★★★★", "★★★★☆", "Hermes"],
        ["插件生态标准化", "★★★☆☆", "★★★★★", "OpenClaw"],
        ["安全纵深防御", "★★★★☆", "★★★★☆", "平手"],
        ["记忆系统深度", "★★★★★", "★★★☆☆", "Hermes"],
        ["运维成熟度", "★★☆☆☆", "★★★★★", "OpenClaw"],
        ["AI 模型兼容性", "★★★★★", "★★★★☆", "Hermes"],
        ["多 Agent 配置", "★★★☆☆", "★★★★★", "OpenClaw"],
        ["配套客户端", "★☆☆☆☆", "★★★★★", "OpenClaw"],
        ["RL 训练环境", "★★★★★", "☆☆☆☆☆", "Hermes"],
    ]

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.3), Inches(1.4), Inches(9.4), Inches(5.5))
    table = table_shape.table
    col_widths = [Inches(2.2), Inches(1.8), Inches(1.8), Inches(1.6)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Microsoft YaHei"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            bg = ROW_ALT if r_idx % 2 == 0 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

            # 根据优势方着色
            winner = row[3]
            text_color = TEXT_DARK
            if c_idx == 3:
                if winner == "Hermes":
                    text_color = ACCENT_HERMES
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xE5)
                elif winner == "OpenClaw":
                    text_color = ACCENT_OPENCLAW
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xE5, 0xFA, 0xF0)

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = text_color
                p.font.name = "Microsoft YaHei"
                if c_idx == 0:
                    p.font.bold = True
                if c_idx >= 1:
                    p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def slide_scenarios(prs):
    """适用场景推荐"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_subtitle(slide, "13 适用场景推荐", "选择的关键在于核心需求")

    # Hermes 场景
    h_card = add_rounded_rect(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(4.5), RGBColor(0xFF, 0xF5, 0xEE))
    tf = set_text(h_card, "选择 Hermes Agent", font_size=16, bold=True, color=ACCENT_HERMES)
    h_card.text_frame.word_wrap = True
    items = [
        "Agent 具备自我学习能力",
        "深度接入中国消息平台",
        "兼容大量非标准 AI 模型",
        "需要 RL 训练环境",
        "个人研究/实验，偏好 Python",
        "ACP 协议集成（IDE 插件）",
    ]
    for item in items:
        add_text_paragraph(tf, f"  •  {item}", font_size=12, color=TEXT_DARK, space_before=Pt(6))

    # OpenClaw 场景
    o_card = add_rounded_rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(4.5), RGBColor(0xEE, 0xFB, 0xF5))
    tf = set_text(o_card, "选择 OpenClaw", font_size=16, bold=True, color=ACCENT_OPENCLAW)
    o_card.text_frame.word_wrap = True
    items = [
        "接入大量国际消息平台",
        "标准化插件生态需求",
        "生产级运维能力",
        "macOS/iOS/Android 配套客户端",
        "团队协作，偏好 TypeScript",
        "内置安全审计命令",
    ]
    for item in items:
        add_text_paragraph(tf, f"  •  {item}", font_size=12, color=TEXT_DARK, space_before=Pt(6))

    # 底部总结
    summary = add_rounded_rect(slide, Inches(0.5), Inches(6.3), Inches(9), Inches(0.7), PRIMARY)
    set_text(summary, '"全能 Agent" (Hermes) vs "万能网关" (OpenClaw) — 取决于核心需求是 "Agent 够强" 还是 "连接够多"', font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_end(prs):
    """结尾"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)

    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = set_text(box, "THANKS", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    sub = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.5))
    set_text(sub, "Hermes Agent vs OpenClaw 架构对比分析", font_size=16, color=RGBColor(0xBB, 0xD5, 0xFF), alignment=PP_ALIGN.CENTER)

    info = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    set_text(info, "分析日期：2026-04-14  |  版本：Hermes v0.9.0 / OpenClaw v2026.4.11", font_size=12, color=RGBColor(0x99, 0xBB, 0xFF), alignment=PP_ALIGN.CENTER)


# ── 主函数 ──

def main():
    """生成 PPT 主函数"""
    output_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "analysis", "hermes-openclaw")
    output_path = os.path.join(output_dir, "Hermes-vs-OpenClaw对比分析.pptx")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides = [
        slide_cover,
        slide_toc,
        slide_positioning,
        slide_tech_stack,
        slide_architecture,
        slide_agent_engine,
        slide_channels,
        slide_tools,
        slide_memory,
        slide_security,
        slide_plugins,
        slide_deploy,
        slide_ai_models,
        slide_multi_agent,
        slide_matrix,
        slide_scenarios,
        slide_end,
    ]

    for fn in slides:
        fn(prs)

    prs.save(output_path)
    print(f"PPT 已生成: {output_path}")
    print(f"共 {len(slides)} 页")


if __name__ == "__main__":
    main()
